#!/usr/bin/env python3
"""
build_and_prepare.py — deterministic publisher helper for pmc-drive-report-publisher.

Purpose: remove model discretion from the "publish to Drive" step. This script
turns approved report content into a REAL Microsoft Office file and prints the
EXACT arguments the agent must pass to the Google Drive `create_file` MCP tool.
The agent must NOT hand-craft those arguments or upload text/html — it copies
what this script emits verbatim.

Why a script can't upload directly: `create_file` is an MCP tool, callable only
by the agent, not from a subprocess. So this script does the deterministic part
(build correct Office bytes, fix the MIME type, force conversion off, base64 it)
and leaves ONLY the single tool call to the agent, with every dangerous knob
pre-filled.

Usage:
  python3 build_and_prepare.py --title "Atlas — Steering Pack — 2026-07-30" \
      --kind docx --content-file /tmp/report_body.md \
      [--parent-id 0AGindXMKcjpZUk9PVA]

--kind is one of: docx | xlsx | pptx
--content-file is UTF-8 text:
    docx  -> Markdown-ish text; headings (#/##), blank-line paragraphs
    xlsx  -> CSV
    pptx  -> one slide per line, "Title | bullet; bullet; bullet"

It writes the Office file next to the content file and prints a JSON block
delimited by <CREATE_FILE_ARGS> ... </CREATE_FILE_ARGS> containing:
  title, parentId, contentMimeType, disableConversionToGoogleType (true),
  base64Content, and the local file path.
"""
import argparse, base64, csv, io, json, os, sys

OFFICE = {
    "docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
    "xlsx": ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
    "pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
}
DEFAULT_PARENT = "0AGindXMKcjpZUk9PVA"  # Atlas Shared Drive root

def build_docx(text, out):
    from docx import Document
    d = Document()
    for block in text.split("\n\n"):
        block = block.strip()
        if not block:
            continue
        if block.startswith("## "):
            d.add_heading(block[3:].strip(), level=2)
        elif block.startswith("# "):
            d.add_heading(block[2:].strip(), level=1)
        else:
            d.add_paragraph(block)
    d.save(out)

def build_xlsx(text, out):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active
    for row in csv.reader(io.StringIO(text)):
        ws.append(row)
    wb.save(out)

def build_pptx(text, out):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        title, _, rest = line.partition("|")
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title.strip()
        body = slide.placeholders[1].text_frame
        bullets = [b.strip() for b in rest.split(";") if b.strip()]
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                body.add_paragraph().text = b
    prs.save(out)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--title", required=True)
    ap.add_argument("--kind", required=True, choices=list(OFFICE))
    ap.add_argument("--content-file", required=True)
    ap.add_argument("--parent-id", default=DEFAULT_PARENT)
    a = ap.parse_args()

    mime, ext = OFFICE[a.kind]
    title = a.title if a.title.endswith(ext) else a.title + ext  # extension is mandatory
    with open(a.content_file, "r", encoding="utf-8") as f:
        text = f.read()

    out = os.path.splitext(a.content_file)[0] + ext
    {"docx": build_docx, "xlsx": build_xlsx, "pptx": build_pptx}[a.kind](text, out)

    with open(out, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()

    args = {
        "title": title,
        "parentId": a.parent_id,
        "contentMimeType": mime,
        "disableConversionToGoogleType": True,   # NON-NEGOTIABLE
        "base64Content": b64,
        "_localFile": out,
    }
    print("<CREATE_FILE_ARGS>")
    print(json.dumps(args))
    print("</CREATE_FILE_ARGS>")
    sys.stderr.write(f"Built {out} ({os.path.getsize(out)} bytes) as {a.kind}. "
                     f"Pass the JSON above to create_file verbatim.\n")

if __name__ == "__main__":
    main()
